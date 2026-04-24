"""
slide_06_company_comparison.py
──────────────────────────────
Slide 6 index + KPI Comparison slides.

After slide 5 (profit pool) the deck gets:
  Slide 6   - KPI Index slide  (python-pptx, inserted by builder.py)
  Slides 7+ - One bar and/or line chart per KPI selection

Scaling note
────────────
ciq_helpers._compute_metrics() stores ratio-based KPIs (EBITDA margin,
ROIC, ROCE, ROA, Debt-to-Equity, % Short-Term Debt) as decimals in wide_df
(e.g. 0.288 = 28.8%).  The KPI registry marks these with "wide_scale": 100
so _build_bar_payload multiplies them back to percentage form before sending
to think-cell.  long_df already holds percentages, so no scaling is needed
for the line chart path.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
from pptx import Presentation

# ── Template paths ─────────────────────────────────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parents[0]

DEFAULT_BAR_TEMPLATE = os.getenv(
    "THINKCELL_TEMPLATE_BAR",
    str(PROJECT_ROOT / "THINKCELL_TEMPLATE_BAR.pptx"),
)
DEFAULT_LINE_TEMPLATE = os.getenv(
    "THINKCELL_TEMPLATE_LINE",
    str(PROJECT_ROOT / "thinkcell_template_growth.pptx"),
)

# ── think-cell element names ───────────────────────────────────────────────────
TC_ELEM_CHART = "BarChart"
TC_ELEM_LINE  = "GrowthChart"
TC_ELEM_TITLE = "ChartTitle"

# ── Static slide titles ────────────────────────────────────────────────────────
TITLE_BAR   = "Company KPI Benchmarking \u2014 Peer Comparison"
TITLE_LINE  = "Company KPI Performance \u2014 Historical Trend"
TITLE_INDEX = "KPI Comparison Overview"

# ── KPI registry ───────────────────────────────────────────────────────────────
# wide_scale: multiplier applied to wide_df values before charting.
#   = 100  -> wide_df stores decimals (0.288), multiply to get % (28.8)
#   absent -> wide_df stores the value as-is (absolute amounts, ratios like x)
KPI_CATEGORIES: Dict[str, List[Dict[str, Any]]] = {
    "Investor Value": [
        {"key": "market_cap",       "label": "Market Capitalization",
         "value_col": "Market capitalization ($ mn)", "yaxis_title": "$ mn",
         "tickformat": ",.0f", "sort": "desc"},
        {"key": "enterprise_value", "label": "Enterprise Value (EV)",
         "value_col": "Enterprise value ($ mn)",      "yaxis_title": "$ mn",
         "tickformat": ",.0f", "sort": "desc"},
    ],
    "Earnings Quality": [
        {"key": "revenue",       "label": "Revenue",
         "value_col": "_Revenue",      "yaxis_title": "$ mn",
         "tickformat": ",.0f", "sort": "desc"},
        {"key": "ebitda",        "label": "EBITDA",
         "value_col": "_EBITDA",       "yaxis_title": "$ mn",
         "tickformat": ",.0f", "sort": "desc"},
        {"key": "ebitda_margin", "label": "EBITDA Margin",
         "value_col": "EBITDA margin", "yaxis_title": "%",
         "tickformat": ".2f",  "sort": "desc", "wide_scale": 100},
        {"key": "yoy_ebitda",    "label": "YoY EBITDA Growth",
         "value_col": "YoY Growth",    "yaxis_title": "%",
         "tickformat": ".2f",  "sort": "desc",
         "level": "company_long", "long_metric": "EBITDA ($ mn)"},
    ],
    "Capital Efficiency": [
        {"key": "roic",           "label": "ROIC",
         "value_col": "ROIC (%)",       "yaxis_title": "%",
         "tickformat": ".2f", "sort": "desc", "wide_scale": 100},
        {"key": "roce",           "label": "ROCE",
         "value_col": "ROCE (%)",       "yaxis_title": "%",
         "tickformat": ".2f", "sort": "desc", "wide_scale": 100},
        {"key": "roa",            "label": "ROA",
         "value_col": "ROA (%)",        "yaxis_title": "%",
         "tickformat": ".2f", "sort": "desc", "wide_scale": 100},
        {"key": "asset_turnover", "label": "Asset Turnover",
         "value_col": "Asset turnover", "yaxis_title": "x",
         "tickformat": ".2f", "sort": "desc"},
    ],
    "Financial Risk": [
        {"key": "net_debt",        "label": "Net Debt",
         "value_col": "Net debt ($ mn)",   "yaxis_title": "$ mn",
         "tickformat": ",.0f", "sort": "asc"},
        {"key": "net_debt_ebitda", "label": "Net Debt / EBITDA",
         "value_col": "Net debt / EBITDA", "yaxis_title": "x",
         "tickformat": ".2f",  "sort": "asc"},
        {"key": "debt_to_equity",  "label": "Debt-to-Equity",
         "value_col": "Debt-to-equity",    "yaxis_title": "%",
         "tickformat": ".2f",  "sort": "asc", "wide_scale": 100},
        {"key": "pct_short_term",  "label": "% Short-Term Debt",
         "value_col": "% short-term debt", "yaxis_title": "%",
         "tickformat": ".2f",  "sort": "asc", "wide_scale": 100},
    ],
    "Cash & Valuation": [
        {"key": "opcf",      "label": "Operating Cash Flow",
         "value_col": "Operating cash flow ($ mn)", "yaxis_title": "$ mn",
         "tickformat": ",.0f", "sort": "desc"},
        {"key": "fcf",       "label": "Free Cash Flow",
         "value_col": "Free cash flow ($ mn)",      "yaxis_title": "$ mn",
         "tickformat": ",.0f", "sort": "desc"},
        {"key": "ev_ebitda", "label": "EV / EBITDA",
         "value_col": "EV / EBITDA",                "yaxis_title": "x",
         "tickformat": ".2f",  "sort": "asc"},
        {"key": "pe",        "label": "P/E",
         "value_col": "P/E",                        "yaxis_title": "x",
         "tickformat": ".2f",  "sort": "asc"},
    ],
    "Workforce Efficiency": [
        {"key": "ebitda_per_fte",  "label": "EBITDA / FTE",
         "value_col": "EBITDA per Employee (USD '000)",
         "yaxis_title": "USD '000", "tickformat": ",.0f", "sort": "desc"},
        {"key": "revenue_per_fte", "label": "Revenue / FTE",
         "value_col": "Revenue per Employee (USD '000)",
         "yaxis_title": "USD '000", "tickformat": ",.0f", "sort": "desc"},
        {"key": "labor_intensity", "label": "Labor Intensity",
         "value_col": "Labor intensity (FTE per $ mn revenue)",
         "yaxis_title": "FTE/$mn", "tickformat": ".2f", "sort": "asc"},
    ],
}

# Flat lookups
KPI_BY_KEY: Dict[str, Dict[str, Any]] = {
    kpi["key"]: kpi
    for kpis in KPI_CATEGORIES.values()
    for kpi in kpis
}
KPI_BY_LABEL: Dict[str, Dict[str, Any]] = {
    kpi["label"]: kpi
    for kpis in KPI_CATEGORIES.values()
    for kpi in kpis
}

CHART_MODE_OPTIONS = ["Point-in-time (bar)", "Time series (line)", "Both"]
CHART_MODE_MAP     = {
    "Point-in-time (bar)": "point_in_time",
    "Time series (line)":  "time_series",
    "Both":                "both",
}
CHART_MODE_REVERSE = {v: k for k, v in CHART_MODE_MAP.items()}


# ── Per-KPI selection dataclass ────────────────────────────────────────────────

@dataclass
class KpiSelection:
    kpi_key:    str
    chart_mode: str   # "point_in_time" | "time_series" | "both"

    @property
    def label(self) -> str:
        return KPI_BY_KEY.get(self.kpi_key, {}).get("label", self.kpi_key)

    @property
    def category(self) -> str:
        for cat, kpis in KPI_CATEGORIES.items():
            if any(k["key"] == self.kpi_key for k in kpis):
                return cat
        return ""

    @property
    def chart_mode_label(self) -> str:
        return CHART_MODE_REVERSE.get(self.chart_mode, self.chart_mode)

    def slides_count(self) -> int:
        return 2 if self.chart_mode == "both" else 1


# ── Master request dataclass ───────────────────────────────────────────────────

@dataclass
class ComparisonSlideRequest:
    base_company:      str
    peer_companies:    List[str]
    kpi_selections:    List[KpiSelection]
    year:              int
    wide_df:           Optional[pd.DataFrame]
    long_df:           Optional[pd.DataFrame]
    bar_template_path:  str = field(default_factory=lambda: DEFAULT_BAR_TEMPLATE)
    line_template_path: str = field(default_factory=lambda: DEFAULT_LINE_TEMPLATE)
    country:            Optional[str] = None
    year_range_start:   int = 2010

    def total_slides(self) -> int:
        return sum(s.slides_count() for s in self.kpi_selections)

    # Back-compat single-KPI accessors
    @property
    def kpi_key(self) -> str:
        return self.kpi_selections[0].kpi_key if self.kpi_selections else ""

    @property
    def kpi_category(self) -> str:
        return self.kpi_selections[0].category if self.kpi_selections else ""

    @property
    def chart_mode(self) -> str:
        return self.kpi_selections[0].chart_mode if self.kpi_selections else "both"


# ── think-cell cell constructors ──────────────────────────────────────────────

def _tc_string(x: Any) -> Dict[str, str]:
    return {"string": str(x)}

def _tc_number(x: Optional[float]) -> Optional[Dict[str, float]]:
    return {"number": float(x)} if x is not None else None


# ── Point-in-time (bar) payload ───────────────────────────────────────────────

def _build_bar_payload(
    req: ComparisonSlideRequest,
    kpi: Dict[str, Any],
    wide_df: pd.DataFrame,
) -> List[dict]:
    """
    Build think-cell bar chart data for one KPI, point-in-time snapshot.

    Applies kpi["wide_scale"] (default 1) to convert decimal-stored metrics
    (e.g. EBITDA margin stored as 0.288) back to percentage form (28.8).
    """
    value_col     = kpi.get("value_col", "")
    all_companies = [req.base_company] + req.peer_companies
    df = wide_df[wide_df["Company"].astype(str).isin(all_companies)].copy()

    # Long-derived KPIs (e.g. YoY growth computed from long_df)
    if kpi.get("level") == "company_long" and req.long_df is not None:
        try:
            # Compute YoY growth from long_df directly
            metric = kpi.get("long_metric", "EBITDA ($ mn)")
            sub = req.long_df[req.long_df["Company"].astype(str).isin([req.base_company] + req.peer_companies)].copy()
            if "Metric" in sub.columns:
                sub = sub[sub["Metric"].astype(str).str.strip() == metric]
            if "Value" in sub.columns and "Year" in sub.columns:
                sub = sub.sort_values("Year")
                sub["_prev"] = sub.groupby("Company")["Value"].shift(1)
                sub["YoY Growth"] = (pd.to_numeric(sub["Value"], errors="coerce") / pd.to_numeric(sub["_prev"], errors="coerce") - 1) * 100
                year_sub = sub[sub["Year"] == req.year][["Company", "YoY Growth"]]
                df = df.merge(year_sub, on="Company", how="left")
                value_col = "YoY Growth"
        except Exception:
            pass

    if value_col not in df.columns:
        return []

    df["_val"] = pd.to_numeric(df[value_col], errors="coerce")
    df = df.dropna(subset=["_val"])

    # ── Scale fix: wide_df stores some % KPIs as decimals (0.288 = 28.8%) ──
    # wide_scale=100 on the KPI definition signals this needs multiplying back.
    wide_scale = kpi.get("wide_scale", 1)
    if wide_scale != 1:
        df["_val"] = df["_val"] * wide_scale

    df = df.sort_values("_val", ascending=(kpi.get("sort", "desc") == "asc"))

    companies = df["Company"].astype(str).tolist()
    values    = df["_val"].tolist()
    if not companies:
        return []

    series_label = kpi.get("label", kpi["key"])
    header_row   = [_tc_string("")] + [_tc_string(c) for c in companies]
    value_row    = [_tc_string(series_label)] + [_tc_number(v) for v in values]

    return [
        {"name": TC_ELEM_CHART, "table": [header_row, value_row]},
        {"name": TC_ELEM_TITLE, "table": [[_tc_string(TITLE_BAR)]]},
    ]


# ── Time-series (line) payload ────────────────────────────────────────────────

def _build_line_payload(
    req: ComparisonSlideRequest,
    kpi: Dict[str, Any],
    long_df: pd.DataFrame,
) -> List[dict]:
    """
    Build think-cell line chart data for one KPI, full time series.
    Uses services.ciq_loader.build_timeseries which calls wide_for_year +
    compute_metrics per year — same pipeline as the KPI page charts.
    """
    from services.ciq_loader import build_timeseries

    all_companies = [req.base_company] + req.peer_companies
    value_col     = kpi.get("value_col", "")
    wide_scale    = kpi.get("wide_scale", 1)

    try:
        available = sorted(long_df["Year"].dropna().astype(int).unique().tolist())
        years = [y for y in available if req.year_range_start <= y <= req.year]
    except Exception:
        years = list(range(req.year_range_start, req.year + 1))

    if not years:
        return []

    value_map: Dict[str, Dict[int, Optional[float]]] = {c: {} for c in all_companies}

    if kpi.get("level") == "company_long":
        # YoY growth — compute from raw metric in long_df
        try:
            metric = kpi.get("long_metric", "EBITDA ($ mn)")
            sub = long_df[long_df["Company"].astype(str).isin(all_companies)].copy()
            if "Metric" in sub.columns:
                sub = sub[sub["Metric"].astype(str).str.strip() == metric]
            sub = sub.sort_values(["Company", "Year"])
            sub["_prev"] = sub.groupby("Company")["Value"].shift(1)
            sub["YoY Growth"] = (
                pd.to_numeric(sub["Value"], errors="coerce")
                / pd.to_numeric(sub["_prev"], errors="coerce") - 1
            ) * 100
            for _, row in sub.iterrows():
                co = str(row.get("Company", ""))
                y  = int(row.get("Year", 0))
                v  = row.get("YoY Growth")
                if co in value_map and y in years:
                    value_map[co][y] = None if pd.isna(v) else float(v)
        except Exception as e:
            print(f"[SLIDE06] YoY growth error: {e}", flush=True)
    else:
        # Use build_timeseries: wide_for_year + compute_metrics per year
        # value_col is a computed column name (e.g. "_EBITDA", "ROIC (%)")
        try:
            ts_rows = build_timeseries(
                long_df=long_df,
                years=years,
                companies=all_companies,
                value_col=value_col,
                country=req.country or "All",
            )
            print(f"[SLIDE06] build_timeseries({value_col}): {len(ts_rows)} rows", flush=True)
            for row in ts_rows:
                co = str(row.get("Company", ""))
                y  = int(row.get("Year", 0))
                v  = row.get("Value")
                if co in value_map and y in years:
                    value_map[co][y] = None if (v is None or pd.isna(v)) else float(v)
        except Exception as e:
            print(f"[SLIDE06] build_timeseries error: {e}", flush=True)

    if not any(value_map[c] for c in all_companies):
        print(f"[SLIDE06] value_map empty for {value_col}", flush=True)
        return []

    # Apply wide_scale if needed (decimal → percentage)
    if wide_scale == 100:
        for c in all_companies:
            for y in list(value_map[c].keys()):
                v = value_map[c][y]
                if v is not None and abs(v) < 2:
                    value_map[c][y] = v * 100

    header_row = [_tc_string("")] + [_tc_string(str(y)) for y in years]
    rows       = [header_row]
    for company in all_companies:
        vals = [_tc_number(value_map[company].get(y)) for y in years]
        rows.append([_tc_string(company)] + vals)

    return [
        {"name": TC_ELEM_LINE,  "table": rows},
        {"name": TC_ELEM_TITLE, "table": [[_tc_string(TITLE_LINE)]]},
    ]


# ── KPI Index slide (python-pptx, styled to match slide_01_exec_summary) ──────

def build_kpi_index_slide(
    prs: "Presentation",  # noqa: F821
    req: ComparisonSlideRequest,
    first_kpi_slide_number: int = 7,
) -> None:
    """
    Append a KPI Index slide to prs (caller moves it into position).

    Styled to match slide_01_exec_summary.py:
      - White background, Calibri font
      - Title 24pt TEXT_DARK + Bain-red rule (identical geometry)
      - Subtitle 13pt TEXT_MID
      - Navy header row + alternating CARD_FILL/white data rows
      - TEXT_DARK bold for KPI name, BAIN_RED bold for slide reference
      - Footer peers line in TEXT_LIGHT
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt

    # ── Colour palette (identical to slide_01) ────────────────────────────────
    BAIN_RED   = RGBColor(225,  28,  42)
    TEXT_DARK  = RGBColor( 17,  24,  39)
    TEXT_MID   = RGBColor( 55,  65,  81)
    TEXT_LIGHT = RGBColor(107, 114, 128)
    CARD_FILL  = RGBColor(250, 250, 252)
    CARD_LINE  = RGBColor(218, 222, 230)
    WHITE      = RGBColor(255, 255, 255)
    BAIN_NAVY  = RGBColor(  0,  50, 109)

    FONT_FACE = "Calibri"

    # ── Geometry (mirrors slide_01 constants) ─────────────────────────────────
    HEADER_X     = 0.31
    HEADER_Y     = 0.38
    HEADER_W     = 12.50
    RED_RULE_Y   = HEADER_Y + 0.57
    RED_RULE_W   = 12.96
    SUBTITLE_Y   = HEADER_Y + 0.62
    CONTENT_TOP  = HEADER_Y + 0.66 + 0.30 + 0.18   # ~1.52

    TABLE_LEFT   = 0.50
    TABLE_W      = 12.33
    HEADER_ROW_H = 0.40
    DATA_ROW_H   = 0.44
    FOOTER_GAP   = 0.18

    COL_FRACS  = [0.06, 0.22, 0.34, 0.24, 0.14]
    COL_WIDTHS = [TABLE_W * f for f in COL_FRACS]
    HEADERS    = ["#", "Category", "KPI", "Chart type", "Slide(s)"]

    FONT_TITLE_PT    = 24
    FONT_SUBTITLE_PT = 13
    FONT_HEADER_PT   = 10
    FONT_BODY_PT     = 10
    FONT_FOOTER_PT   =  9

    def _in(val: float) -> int:
        return int(Inches(val))

    # ── XML bullet-suppression helpers (same as slide_01) ─────────────────────
    def _force_no_bullets(para):
        try:
            pPr = para._p.get_or_add_pPr()
            for tag in ("a:buChar", "a:buAutoNum", "a:buBlip", "a:buFont",
                        "a:buSzPct", "a:buSzPts", "a:buClr", "a:buNone"):
                for el in pPr.findall(tag):
                    pPr.remove(el)
            pPr.append(OxmlElement("a:buNone"))
            para.level = 0
        except Exception:
            pass

    def _nuke_bullets(tf):
        try:
            txBody = tf._txBody
            for old in txBody.findall(qn("a:lstStyle")):
                txBody.remove(old)
            lstStyle = OxmlElement("a:lstStyle")
            for lvl in range(1, 10):
                lvlPPr = OxmlElement(f"a:lvl{lvl}pPr")
                lvlPPr.append(OxmlElement("a:buNone"))
                lstStyle.append(lvlPPr)
            bodyPr = txBody.find(qn("a:bodyPr"))
            if bodyPr is not None:
                bodyPr.addnext(lstStyle)
            else:
                txBody.insert(0, lstStyle)
        except Exception:
            pass

    def _set_spacing(para, before=0, after=0):
        try:
            pPr = para._p.get_or_add_pPr()
            for attr, val in (("a:spcBef", before), ("a:spcAft", after)):
                for old in pPr.findall(attr):
                    pPr.remove(old)
                spc    = OxmlElement(attr)
                spcPts = OxmlElement("a:spcPts")
                spcPts.set("val", str(int(val * 100)))
                spc.append(spcPts)
                pPr.append(spc)
        except Exception:
            pass

    # ── Blank layout lookup (same logic as slide_01._get_blank_layout) ────────
    blank_layout = None
    for _l in prs.slide_layouts:
        try:
            if (getattr(_l, "name", "") or "").strip().casefold() == "blank":
                blank_layout = _l
                break
        except Exception:
            pass
    if blank_layout is None:
        for _l in prs.slide_layouts:
            try:
                if len(getattr(_l, "placeholders", [])) == 0:
                    blank_layout = _l
                    break
            except Exception:
                pass
    if blank_layout is None:
        blank_layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]

    slide = prs.slides.add_slide(blank_layout)

    # White background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Remove inherited placeholders
    for shp in [s for s in slide.shapes if getattr(s, "is_placeholder", False)]:
        try:
            slide.shapes._spTree.remove(shp._element)
        except Exception:
            pass

    # ── Helper: single-run textbox ────────────────────────────────────────────
    def _tb(l, t, w, h, text, pt, bold=False, italic=False,
            color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(_in(l), _in(t), _in(w), _in(h))
        tf  = txb.text_frame
        tf.word_wrap = wrap
        _nuke_bullets(tf)
        p = tf.paragraphs[0]
        _force_no_bullets(p)
        _set_spacing(p)
        p.alignment = align
        r = p.add_run()
        r.text           = text
        r.font.size      = Pt(pt)
        r.font.bold      = bold
        r.font.italic    = italic
        r.font.color.rgb = color
        r.font.name      = FONT_FACE
        return txb

    # ── Helper: filled rectangle ───────────────────────────────────────────────
    def _rect(l, t, w, h, fill, line_color=None, line_pt=0.75, rounded=False):
        st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(st, _in(l), _in(t), _in(w), _in(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line_color:
            shp.line.color.rgb = line_color
            shp.line.width     = Pt(line_pt)
        else:
            shp.line.fill.background()
        return shp

    # ── Title ──────────────────────────────────────────────────────────────────
    _tb(HEADER_X, HEADER_Y, HEADER_W, 0.60,
        TITLE_INDEX, FONT_TITLE_PT, bold=True, color=TEXT_DARK)

    # ── Red rule ───────────────────────────────────────────────────────────────
    _rect(0, RED_RULE_Y, RED_RULE_W, 0.02, fill=BAIN_RED)

    # ── Subtitle ───────────────────────────────────────────────────────────────
    n_peers  = len(req.peer_companies)
    peer_txt = f"{n_peers} peer{'s' if n_peers != 1 else ''}"
    sub_text = (
        f"{req.base_company} vs {peer_txt}"
        f"  \u2502  FY{req.year}"
        f"  \u2502  {req.year_range_start}\u2013{req.year} time series"
    )
    if req.country:
        sub_text = f"{req.country}  \u2502  " + sub_text
    _tb(HEADER_X, SUBTITLE_Y, HEADER_W, 0.32,
        sub_text, FONT_SUBTITLE_PT, color=TEXT_MID)

    # ── Compute cumulative slide numbers ───────────────────────────────────────
    rows_data: List[Tuple[int, str, str, str, str]] = []
    slide_num = first_kpi_slide_number
    for sel in req.kpi_selections:
        n = sel.slides_count()
        slide_ref = str(slide_num) if n == 1 else f"{slide_num}\u2013{slide_num + n - 1}"
        rows_data.append((len(rows_data) + 1, sel.category, sel.label,
                          sel.chart_mode_label, slide_ref))
        slide_num += n

    # ── Table: header row ─────────────────────────────────────────────────────
    PAD_X   = 0.14
    PAD_Y_H = 0.10
    PAD_Y_D = 0.12

    x = TABLE_LEFT
    for col_w, hdr in zip(COL_WIDTHS, HEADERS):
        _rect(x, CONTENT_TOP, col_w, HEADER_ROW_H,
              fill=BAIN_NAVY, line_color=WHITE, line_pt=0.5)
        _tb(x + PAD_X, CONTENT_TOP + PAD_Y_H,
            col_w - PAD_X - 0.04, HEADER_ROW_H - PAD_Y_H,
            hdr, FONT_HEADER_PT, bold=True, color=WHITE, wrap=False)
        x += col_w

    # ── Table: data rows ──────────────────────────────────────────────────────
    for i, (num, cat, kpi_lbl, mode_lbl, sref) in enumerate(rows_data):
        row_top = CONTENT_TOP + HEADER_ROW_H + i * DATA_ROW_H
        bg      = CARD_FILL if i % 2 == 0 else WHITE
        values  = [str(num), cat, kpi_lbl, mode_lbl, sref]

        x = TABLE_LEFT
        for col_w, val in zip(COL_WIDTHS, values):
            _rect(x, row_top, col_w, DATA_ROW_H,
                  fill=bg, line_color=CARD_LINE, line_pt=0.5)

            if val == kpi_lbl:
                txt_color, txt_bold = TEXT_DARK, True
            elif val == sref:
                txt_color, txt_bold = BAIN_RED, True
            else:
                txt_color, txt_bold = TEXT_MID, False

            _tb(x + PAD_X, row_top + PAD_Y_D,
                col_w - PAD_X - 0.04, DATA_ROW_H - PAD_Y_D,
                val, FONT_BODY_PT, bold=txt_bold, color=txt_color, wrap=False)
            x += col_w

    # ── Footer: peers ─────────────────────────────────────────────────────────
    footer_top = CONTENT_TOP + HEADER_ROW_H + len(rows_data) * DATA_ROW_H + FOOTER_GAP
    peers_str  = ", ".join(req.peer_companies[:6])
    if len(req.peer_companies) > 6:
        peers_str += f" +{len(req.peer_companies) - 6} more"
    _tb(TABLE_LEFT, footer_top, TABLE_W, 0.28,
        f"Peers: {peers_str}", FONT_FOOTER_PT, color=TEXT_LIGHT)


# ── Public API: think-cell payloads ───────────────────────────────────────────

def build_slide_06_tc_payloads(req: ComparisonSlideRequest) -> List[Dict[str, Any]]:
    """
    Return flat list of {template, data} dicts for ALL selected KPIs.
    Does NOT include the index slide — that is added by builder.py via
    build_kpi_index_slide().

    For each KpiSelection:
      point_in_time -> 1 bar-chart item
      time_series   -> 1 line-chart item
      both          -> bar item + line item (2 slides)
    """
    results: List[Dict[str, Any]] = []

    for sel in req.kpi_selections:
        kpi = KPI_BY_KEY.get(sel.kpi_key)
        print(f"[SLIDE06] sel={sel.kpi_key} mode={sel.chart_mode} kpi_found={kpi is not None}", flush=True)
        if kpi is None:
            continue

        if sel.chart_mode in ("point_in_time", "both"):
            wide_ok = req.wide_df is not None and not req.wide_df.empty
            print(f"[SLIDE06] bar: wide_ok={wide_ok}", flush=True)
            if wide_ok:
                data = _build_bar_payload(req, kpi, req.wide_df)
                print(f"[SLIDE06] bar data items={len(data)}", flush=True)
                if data:
                    results.append({"template": req.bar_template_path, "data": data})

        if sel.chart_mode in ("time_series", "both"):
            long_ok = req.long_df is not None and not req.long_df.empty
            print(f"[SLIDE06] line: long_ok={long_ok}", flush=True)
            if long_ok:
                data = _build_line_payload(req, kpi, req.long_df)
                print(f"[SLIDE06] line data items={len(data)}", flush=True)
                if data:
                    results.append({"template": req.line_template_path, "data": data})

    print(f"[SLIDE06] total payloads={len(results)}", flush=True)
    return results