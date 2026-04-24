# slide_04_cement_sales_growth.py

from __future__ import annotations

import io
from typing import Any, Dict, List, Optional

import matplotlib.pyplot as plt
from matplotlib.ticker import PercentFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from deck_builder.data_adapters import GrowthSlideContent


# ── Colours ───────────────────────────────────────────────────────────────────
COLOR_HIST     = "black"
COLOR_FORECAST = "#d60000"
COLOR_COVID    = "gray"
COLOR_TITLE    = RGBColor(17, 24, 39)
COLOR_SUBTITLE = RGBColor(60, 60, 60)

# ── Typography ────────────────────────────────────────────────────────────────
FONT_PT_TITLE = 26
FONT_PT_BODY  = 10
FONT_PT_AXIS  = 11

# ── Chart style ───────────────────────────────────────────────────────────────
CHART_FIGSIZE  = (12.5, 5.2)
CHART_DPI      = 220
LINE_WIDTH     = 2.2
MARKER_SIZE    = 4
COVID_YEARS    = (2019, 2020)
COVID_LABEL_X  = 2019.5
DEFAULT_CUTOFF = 2024

# ── Slide geometry (inches) ───────────────────────────────────────────────────
TITLE_X, TITLE_Y, TITLE_W, TITLE_H = 0.6, 0.35, 12.2, 0.9
HEAD_X,  HEAD_Y,  HEAD_W,  HEAD_H  = 0.6, 1.25, 12.2, 0.6
CHART_X, CHART_Y, CHART_W, CHART_H = 0.75, 1.35, 12.2, 5.2


# ── Series normalisation ──────────────────────────────────────────────────────

def _as_list(x: Any) -> List[Any]:
    if x is None:
        return []
    return list(x) if isinstance(x, (list, tuple)) else [x]


def _normalize_series(series: Any) -> Dict[str, Any]:
    """
    Normalise a growth series into a plain dict with keys:
    years, yoy, revenue, cutoff_year, title, subtitle.

    Accepts either a dict or a dataclass/object with matching attributes.
    """
    if series is None:
        return {"years": [], "yoy": [], "revenue": [], "cutoff_year": DEFAULT_CUTOFF}

    def _parse_years(raw):
        return [int(y) for y in _as_list(raw) if str(y).strip()]

    def _parse_floats(raw):
        return [None if v is None else float(v) for v in _as_list(raw)]

    if isinstance(series, dict):
        return {
            "years":       _parse_years(series.get("years")),
            "yoy":         _parse_floats(series.get("yoy")),
            "revenue":     _parse_floats(series.get("revenue")),
            "cutoff_year": int(series.get("cutoff_year") or DEFAULT_CUTOFF),
        }

    if hasattr(series, "years") and hasattr(series, "yoy"):
        return {
            "years":       _parse_years(getattr(series, "years", [])),
            "yoy":         _parse_floats(getattr(series, "yoy", [])),
            "revenue":     _parse_floats(getattr(series, "revenue", [])),
            "cutoff_year": int(getattr(series, "cutoff_year", None) or DEFAULT_CUTOFF),
            "title":       getattr(series, "title", None),
            "subtitle":    getattr(series, "subtitle", None),
        }

    raise TypeError(f"Expected growth series as dict or object with years+yoy, got {type(series)}")


# ── Chart rendering ───────────────────────────────────────────────────────────

def _split_hist_forecast(
    years: List[int],
    yoy: List[Optional[float]],
    cutoff: int,
) -> tuple:
    """
    Split year/yoy pairs into historical and forecast lists.
    The cutoff year appears in both so the two lines connect.
    """
    hist_x, hist_y, fc_x, fc_y = [], [], [], []
    for y, v in zip(years, yoy):
        if y <= cutoff:
            hist_x.append(y)
            hist_y.append(v)
        if y >= cutoff:
            fc_x.append(y)
            fc_y.append(v)
    return hist_x, hist_y, fc_x, fc_y


def _plot_growth_chart(series: Any) -> bytes:
    """
    Render a slide-ready YoY growth line chart and return PNG bytes.

    - Historical (≤ cutoff_year): solid black line
    - Forecast   (≥ cutoff_year): dashed red line
    - Lines share the cutoff year point so they connect visually
    - COVID-19 region marked with two dashed vertical lines
    """
    s      = _normalize_series(series)
    years  = s["years"]
    yoy    = s["yoy"]
    cutoff = s["cutoff_year"]

    if not years or len(years) != len(yoy):
        years, yoy = [], []

    hist_x, hist_y, fc_x, fc_y = _split_hist_forecast(years, yoy, cutoff)

    fig, ax = plt.subplots(figsize=CHART_FIGSIZE, dpi=CHART_DPI)

    if hist_x:
        ax.plot(hist_x, hist_y, color=COLOR_HIST, linewidth=LINE_WIDTH,
                marker="o", markersize=MARKER_SIZE,
                label=f"Historical ({hist_x[0]}–{cutoff})")

    if fc_x:
        ax.plot(fc_x, fc_y, color=COLOR_FORECAST, linewidth=LINE_WIDTH,
                linestyle="--", marker="o", markersize=MARKER_SIZE,
                label=f"Forecast ({cutoff}–{fc_x[-1]})")

    for vline in COVID_YEARS:
        ax.axvline(x=vline, color=COLOR_COVID, linestyle=(0, (4, 4)), linewidth=1.2, alpha=0.8)

    if years:
        ax.text(COVID_LABEL_X, ax.get_ylim()[1] * 0.98, "COVID-19",
                ha="center", va="top", fontsize=FONT_PT_BODY, color=COLOR_COVID)

    ax.set_xlabel("Year",               fontsize=FONT_PT_AXIS)
    ax.set_ylabel("YoY revenue growth", fontsize=FONT_PT_AXIS)
    ax.yaxis.set_major_formatter(PercentFormatter(xmax=1.0, decimals=0))
    ax.grid(True,  axis="y", alpha=0.25)
    ax.grid(False, axis="x")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend(loc="upper left", frameon=False, fontsize=FONT_PT_BODY, ncol=2)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True)
    plt.close(fig)
    return buf.getvalue()


# ── Slide builder ─────────────────────────────────────────────────────────────

def add_slide_04_cement_sales_growth(
    prs: Presentation,
    theme: Any,
    content: GrowthSlideContent,
) -> None:
    """Slide 4 – Cement / Concrete / Lime YoY Growth View."""
    slide = prs.slides.add_slide(prs.slide_layouts[8])

    # ── Title block ───────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(
        Inches(TITLE_X), Inches(TITLE_Y), Inches(TITLE_W), Inches(TITLE_H)
    )
    tf = title_box.text_frame
    tf.clear()

    title_run            = tf.paragraphs[0].add_run()
    title_run.text       = "Cement / concrete / lime — Growth view"
    title_run.font.size  = Pt(FONT_PT_TITLE)
    title_run.font.bold  = True
    title_run.font.color.rgb = COLOR_TITLE

    if content.country:
        sub_label = (
            f"{content.country} (Region: {content.region})"
            if content.region else content.country
        )
        p2                = tf.add_paragraph()
        p2.text           = sub_label
        p2.level          = 1
        p2.font.size      = Pt(FONT_PT_BODY)
        p2.font.color.rgb = COLOR_SUBTITLE

    # ── Headline ──────────────────────────────────────────────────
    if content.headline:
        hbox = slide.shapes.add_textbox(
            Inches(HEAD_X), Inches(HEAD_Y), Inches(HEAD_W), Inches(HEAD_H)
        )
        htf = hbox.text_frame
        htf.clear()
        hp                = htf.paragraphs[0]
        hp.text           = content.headline
        hp.font.size      = Pt(FONT_PT_BODY)
        hp.font.color.rgb = COLOR_SUBTITLE

    # ── Chart image ───────────────────────────────────────────────
    slide.shapes.add_picture(
        io.BytesIO(_plot_growth_chart(content)),
        Inches(CHART_X), Inches(CHART_Y),
        width=Inches(CHART_W), height=Inches(CHART_H),
    )