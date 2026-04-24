# slide_01_exec_summary.py

from __future__ import annotations

from typing import Any, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


# ── Colour palette ────────────────────────────────────────────────────────────
BAIN_RED    = RGBColor(225,  28,  42)
TEXT_DARK   = RGBColor( 17,  24,  39)
TEXT_MID    = RGBColor( 55,  65,  81)
TEXT_LIGHT  = RGBColor(107, 114, 128)
CARD_FILL   = RGBColor(250, 250, 252)
CARD_LINE   = RGBColor(218, 222, 230)

# ── Typography ────────────────────────────────────────────────────────────────
FONT_TITLE          = "Calibri"
FONT_PT_SLIDE_TITLE = 26
FONT_PT_SUBTITLE    = 13
FONT_PT_SECTION     = 14
FONT_PT_HEADLINE    = 12
FONT_PT_BODY        = 10.0
FONT_PT_BODY_MIN    = 9.0

# ── Content limits ────────────────────────────────────────────────────────────
MAX_BULLETS      = 3
MAX_BULLET_CHARS = 135

# ── Slide geometry (inches, 13.33 × 7.5 widescreen) ─────────────────────────
SLIDE_W          = 13.33
SLIDE_H          = 7.5
HEADER_X         = 0.15
HEADER_Y         = 0.38
HEADER_W         = 12.33
HEADER_INNER_X   = HEADER_X + 0.16
HEADER_INNER_W   = HEADER_W - 0.16
RED_RULE_Y_DELTA = 0.57   # below header_y
RED_RULE_W       = 12.96
CARD_X           = 0.50
CARD_W           = 12.33
CARD_GAP         = 0.08
CARD_BOTTOM_PAD  = 0.20
CARD_HEIGHT_SHRINK = 0.95
CARD_PAD_X       = 0.28
CARD_PAD_TOP     = 0.08


# ── Slide layout helpers ──────────────────────────────────────────────────────

def _get_blank_layout(prs: Presentation):
    """Return the blank slide layout, falling back gracefully."""
    for layout in prs.slide_layouts:
        try:
            if (getattr(layout, "name", "") or "").strip().casefold() == "blank":
                return layout
        except Exception:
            pass
    for layout in prs.slide_layouts:
        try:
            if len(getattr(layout, "placeholders", [])) == 0:
                return layout
        except Exception:
            pass
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def _remove_shape(slide, shape) -> None:
    try:
        slide.shapes._spTree.remove(shape._element)
    except Exception:
        try:
            if shape.has_text_frame:
                shape.text_frame.clear()
        except Exception:
            pass


def _clear_all_placeholders(slide) -> None:
    to_remove = [s for s in slide.shapes if getattr(s, "is_placeholder", False)]
    for shp in to_remove:
        _remove_shape(slide, shp)


# ── Content extraction helpers ────────────────────────────────────────────────

def _safe_get_blocks(content: Any) -> List[Tuple[str, str, List[str]]]:
    """Normalise arbitrary content objects into (title, headline, bullets) tuples."""
    if content is None:
        return []
    blocks = getattr(content, "blocks", None)
    if isinstance(blocks, list) and blocks:
        out: List[Tuple[str, str, List[str]]] = []
        for blk in blocks:
            title       = str(getattr(blk, "title",    "") or "Section").strip()
            headline    = str(getattr(blk, "headline", "") or "").strip()
            bullets_raw = getattr(blk, "bullets", None)
            if isinstance(bullets_raw, list):
                bullets = [str(b).strip() for b in bullets_raw if str(b).strip()]
            elif bullets_raw is not None and str(bullets_raw).strip():
                bullets = [str(bullets_raw).strip()]
            else:
                bullets = []
            out.append((title, headline, bullets))
        return out
    txt = str(getattr(content, "text", "") or content or "").strip()
    return [("Summary", "", [txt])] if txt else []


def _strip_leading_bullet(s: str) -> str:
    t = (s or "").strip()
    for pref in ("•", "·", "◦", "▪", "–", "-", "●", "‣", "*"):
        if t.startswith(pref):
            return t[len(pref):].strip()
    return t


def _truncate(s: str, max_chars: int) -> str:
    s = (s or "").strip()
    if len(s) <= max_chars:
        return s
    cut = s[:max_chars - 1].rsplit(" ", 1)[0]
    if len(cut) < max_chars * 0.6:
        cut = s[:max_chars - 1]
    return cut.rstrip() + "…"


# ── XML / text-frame helpers ──────────────────────────────────────────────────

def _force_no_bullets(paragraph) -> None:
    """Suppress inherited bullets at the paragraph level."""
    try:
        pPr = paragraph._p.get_or_add_pPr()
        for tag in (
            "a:buChar", "a:buAutoNum", "a:buBlip",
            "a:buFont",  "a:buSzPct",  "a:buSzPts",
            "a:buClr",   "a:buNone",
        ):
            for el in pPr.findall(tag):
                pPr.remove(el)
        pPr.append(OxmlElement("a:buNone"))
        paragraph.level = 0
    except Exception:
        pass


def _nuke_textbox_bullets(text_frame) -> None:
    """
    Override bullet inheritance at the txBody <a:lstStyle> level so that
    slide-master bullets cannot cascade into text boxes.
    """
    try:
        txBody = text_frame._txBody
        for old in txBody.findall(qn("a:lstStyle")):
            txBody.remove(old)
        lstStyle = OxmlElement("a:lstStyle")
        for lvl in range(1, 10):
            lvlPPr = OxmlElement(f"a:lvl{lvl}pPr")
            lvlPPr.append(OxmlElement("a:buNone"))
            lstStyle.append(lvlPPr)
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.addnext(lstStyle)
        else:
            txBody.insert(0, lstStyle)
    except Exception:
        pass


def _set_paragraph_spacing(paragraph, space_before_pt: float = 0, space_after_pt: float = 0) -> None:
    """Set paragraph spacing via XML to avoid line-spacing quirks."""
    try:
        pPr = paragraph._p.get_or_add_pPr()
        for attr, val in (("a:spcBef", space_before_pt), ("a:spcAft", space_after_pt)):
            for old in pPr.findall(attr):
                pPr.remove(old)
            spc     = OxmlElement(attr)
            spcPts  = OxmlElement("a:spcPts")
            spcPts.set("val", str(int(val * 100)))
            spc.append(spcPts)
            pPr.append(spc)
    except Exception:
        pass


def _add_text_run(
    text_frame,
    text: str,
    *,
    font_pt: float,
    color: RGBColor,
    bold: bool = False,
    italic: bool = False,
    word_wrap: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Populate a fresh text frame with a single styled run."""
    text_frame.clear()
    _nuke_textbox_bullets(text_frame)
    text_frame.word_wrap = word_wrap
    p = text_frame.paragraphs[0]
    _force_no_bullets(p)
    _set_paragraph_spacing(p)
    p.alignment = align
    r = p.add_run()
    r.text           = text
    r.font.size      = Pt(font_pt)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    r.font.name      = FONT_TITLE


def _add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_pt=0.75, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width     = Pt(line_pt)
    else:
        shp.line.fill.background()
    return shp


# ── Layout maths ─────────────────────────────────────────────────────────────

def _estimate_lines(text: str, chars_per_line: int) -> int:
    if not text:
        return 0
    return max(1, (len(text.strip()) + chars_per_line - 1) // chars_per_line)


def _fit_bullets_to_box(
    *,
    bullets: List[str],
    box_w_in: float,
    box_h_in: float,
    start_font_pt: float = FONT_PT_BODY,
) -> Tuple[List[str], float]:
    bulls = [_truncate(b, MAX_BULLET_CHARS) for b in bullets if b.strip()][:MAX_BULLETS]
    if not bulls:
        bulls = ["(No content available)"]

    chars_per_line = max(28, int(box_w_in * 15))
    total_lines    = sum(_estimate_lines(b, chars_per_line) for b in bulls)

    def max_lines_for(pt: float) -> int:
        return max(1, int(box_h_in / ((1.18 * pt) / 72.0)))

    font_pt = float(start_font_pt)
    while total_lines > max_lines_for(font_pt) and font_pt > FONT_PT_BODY_MIN:
        font_pt -= 0.5

    if total_lines > max_lines_for(font_pt):
        bulls[-1] = _truncate(bulls[-1], 95)

    return bulls, font_pt


# ── Slide building blocks ─────────────────────────────────────────────────────

def _add_title_area(slide, *, top_line: str, subtitle: Optional[str] = None) -> float:
    """
    Render the slide header (title + optional subtitle + red underline rule).
    Returns the Y coordinate at which card content may begin.
    """
    title_clean    = _strip_leading_bullet(top_line or "Executive Summary")
    subtitle_clean = _strip_leading_bullet(subtitle or "")

    # Main title
    tbox = slide.shapes.add_textbox(
        Inches(HEADER_INNER_X), Inches(HEADER_Y),
        Inches(HEADER_INNER_W), Inches(0.60),
    )
    _add_text_run(tbox.text_frame, title_clean, font_pt=FONT_PT_SLIDE_TITLE, color=TEXT_DARK)

    # Red rule
    _add_rect(slide, x=0, y=HEADER_Y + RED_RULE_Y_DELTA, w=RED_RULE_W, h=0.02, fill_rgb=BAIN_RED)

    if subtitle_clean:
        sbox = slide.shapes.add_textbox(
            Inches(HEADER_INNER_X), Inches(HEADER_Y + 0.62),
            Inches(HEADER_INNER_W), Inches(0.30),
        )
        _add_text_run(sbox.text_frame, subtitle_clean, font_pt=FONT_PT_SUBTITLE, color=TEXT_MID)

    return HEADER_Y + 0.66 + 0.30 + 0.10


def _add_block_card(
    slide,
    *,
    title: str,
    headline: str,
    bullets: List[str],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    """Render one content card: section title → italic headline → bullet body."""
    # Background
    _add_rect(slide, x, y, w, h, fill_rgb=CARD_FILL, line_rgb=CARD_LINE, line_pt=0.75, radius=True)

    text_x = x + CARD_PAD_X
    text_w = w - CARD_PAD_X - 0.18

    # Section title
    tbox = slide.shapes.add_textbox(
        Inches(text_x), Inches(y + CARD_PAD_TOP),
        Inches(text_w), Inches(0.30),
    )
    _add_text_run(tbox.text_frame, _strip_leading_bullet(title), font_pt=FONT_PT_SECTION, color=TEXT_DARK, bold=True)

    # Headline (italic, muted)
    headline_h = 0.0
    hh = _strip_leading_bullet(_truncate(headline or "", 130))
    if hh:
        hbox = slide.shapes.add_textbox(
            Inches(text_x), Inches(y + CARD_PAD_TOP + 0.31),
            Inches(text_w), Inches(0.28),
        )
        _add_text_run(hbox.text_frame, hh, font_pt=FONT_PT_HEADLINE, color=TEXT_LIGHT, italic=True)
        headline_h = 0.30

    # Bullet body
    body_y = y + CARD_PAD_TOP + 0.31 + headline_h + 0.08
    body_h = h - (body_y - y) - 0.14
    bulls_fit, font_pt = _fit_bullets_to_box(bullets=bullets, box_w_in=text_w, box_h_in=body_h)

    bbox = slide.shapes.add_textbox(Inches(text_x), Inches(body_y), Inches(text_w), Inches(body_h))
    tfb  = bbox.text_frame
    tfb.clear()
    _nuke_textbox_bullets(tfb)
    tfb.word_wrap = True

    for i, b in enumerate(bulls_fit):
        para = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        _force_no_bullets(para)
        _set_paragraph_spacing(para, space_before_pt=0.5, space_after_pt=1.5)
        para.level = 0

        for run_text, run_color in (("• ", TEXT_DARK), (b, TEXT_MID)):
            r            = para.add_run()
            r.text       = run_text
            r.font.size  = Pt(font_pt)
            r.font.color.rgb = run_color
            r.font.bold  = False
            r.font.name  = FONT_TITLE


# ── Public API ────────────────────────────────────────────────────────────────

_FALLBACK_BLOCK = (
    "Data not available",
    "",
    ["Demand: Data not available.", "Policy & funding: Data not available.", "Supply & costs: Data not available."],
)


def add_slide_01_exec_summary(
    prs: Presentation,
    theme: Any,
    content: Any,
    subtitle: Optional[str] = None,
    *,
    top_line: Optional[str] = None,
) -> None:
    """
    Slide 1 – Executive Summary.

    Layout: slide title + red underline rule, then 3 equal-height content cards
    stacked vertically. Each card contains a section title, italic headline, and
    up to 3 body bullets.
    """
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _clear_all_placeholders(slide)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    content_start_y = _add_title_area(
        slide,
        top_line=top_line or "Executive Summary",
        subtitle=subtitle,
    )

    blocks = (_safe_get_blocks(content) or [_FALLBACK_BLOCK])[:3]
    while len(blocks) < 3:
        blocks.append(_FALLBACK_BLOCK)

    available_h = SLIDE_H - content_start_y - CARD_BOTTOM_PAD
    h = ((available_h - 2 * CARD_GAP) / 3) * CARD_HEIGHT_SHRINK

    for idx, (t, hh, bulls) in enumerate(blocks):
        _add_block_card(
            slide,
            title=t,
            headline=hh,
            bullets=(bulls or [])[:MAX_BULLETS],
            x=CARD_X,
            y=content_start_y + idx * (h + CARD_GAP),
            w=CARD_W,
            h=h,
        )