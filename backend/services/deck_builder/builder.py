"""
services/deck_builder/builder.py

Port of Streamlit builder.py for FastAPI.

Slide assembly order:
  1  Executive Summary         (python-pptx, prepended after TC call)
  2  Construction YoY growth   (think-cell)
  3  Building products growth  (think-cell)
  4  Cement / lime growth      (think-cell)
  5  Profit pool Mekko         (think-cell, optional)
  6  KPI Index                 (python-pptx, only when KPI slides present)
  7+ KPI comparison slides     (think-cell)
"""
from __future__ import annotations

import io
import json
import re as _re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

import pandas as pd
import requests
from pptx import Presentation
from pptx.oxml.ns import qn

from core.config import settings
from services.deck_builder import data_adapters as da
from services.deck_builder.ppt_theme import default_theme
from services.deck_builder.slides.slide_01_exec_summary import add_slide_01_exec_summary
from services.deck_builder.slides.slide_06_company_comparison import (
    ComparisonSlideRequest,
    build_kpi_index_slide,
    build_slide_06_tc_payloads,
)

# ── Constants ─────────────────────────────────────────────────────────────────
TC_ELEM_CHART          = "GrowthChart"
TC_ELEM_TITLE          = "ChartTitle"
TC_ELEM_CATEGORY_CHART = "CategoryChart"
TC_ELEM_CATEGORY_TITLE = "ChartTitle"

COL_COUNTRY   = "Headquarters - Country/Region"
COL_FINAL_TAG = "Final tagging"
COL_SIC       = "SIC Codes"

CEMENT_PARENT    = "Building Materials (Except Cement)"
CATEGORY_RENAMES = {
    "Materials & Components": "Building Materials (Except Cement)",
    "Epc & Design":           "EPC & Design",
}
DEFAULT_PROFIT_POOL_YEAR = 2024

GROWTH_START_YEAR = 2009
GROWTH_END_YEAR   = 2029
HIST_START        = 2010
HIST_END          = 2024
FC_START          = 2025
FC_END            = 2029


# ── DeckRequest ───────────────────────────────────────────────────────────────

@dataclass(frozen=True)
class DeckRequest:
    country:               Optional[str]           = None
    company:               Optional[str]           = None
    year:                  Optional[int]           = None
    df_flat:               Optional[pd.DataFrame]  = field(default=None, compare=False)
    df_profit_pool:        Optional[pd.DataFrame]  = field(default=None, compare=False)
    comparison_request:    Optional[ComparisonSlideRequest] = field(default=None, compare=False)
    tc_server_url:         str = ""
    tc_growth_template_s2: str = ""
    tc_growth_template_s3: str = ""
    tc_growth_template_s4: str = ""
    tc_category_template:  str = ""


# ── think-cell helpers ────────────────────────────────────────────────────────

def _normalize_tc_url(url: str) -> str:
    u = (url or "").strip() or str(settings.THINKCELL_SERVER_URL)
    return u if u.endswith("/") else u + "/"


def _template_to_absolute_url(ref: str) -> str:
    ref = (ref or "").strip()
    if not ref:
        return ref
    if ref.lower().startswith(("http://", "https://", "file:///")):
        return ref
    p = Path(ref)
    abs_path = p if p.is_absolute() else Path.cwd() / p
    return abs_path.resolve().as_uri()


def _tc_string(x: Any) -> Dict[str, str]:
    return {"string": str(x)}


def _tc_year(y: int) -> Dict[str, str]:
    return {"string": str(int(y))}


def _tc_number(x: Optional[float]) -> Optional[Dict[str, float]]:
    return {"number": float(x)} if x is not None else None


def _tc_percent(x: Optional[float]) -> Optional[Dict[str, float]]:
    return {"percentage": float(x)} if x is not None else None


def _build_multi_ppttc_payload(items: List[Dict[str, Any]]) -> str:
    return json.dumps(
        [{"template": _template_to_absolute_url(str(it["template"])), "data": it["data"]}
         for it in items],
        ensure_ascii=False,
    )


def _post_to_tcserver(tcserver_url: str, ppttc_json: str, timeout_sec: int = 180) -> bytes:
    url = _normalize_tc_url(tcserver_url)
    headers = {
        "Content-Type": "application/vnd.think-cell.ppttc+json",
        "Accept":       "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    r = requests.post(url, data=ppttc_json.encode("utf-8"), headers=headers, timeout=timeout_sec)
    if not r.ok:
        detail = (r.text or "").strip() or "(empty response body)"
        raise RuntimeError(f"think-cell Server HTTP {r.status_code}\nURL: {url}\n{detail}")
    return r.content


# ── Post-processing ───────────────────────────────────────────────────────────

def _fix_axis_font_sizes(prs: Presentation, target_pt: float = 10.0) -> None:
    sz_val = str(int(target_pt * 100))
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                for ax_tag in ("c:catAx", "c:valAx"):
                    for ax in shape._element.iter(qn(ax_tag)):
                        for rPr    in ax.iter(qn("a:rPr")):    rPr.set("sz", sz_val)
                        for defRPr in ax.iter(qn("a:defRPr")): defRPr.set("sz", sz_val)
            except Exception:
                pass


def _move_slide(prs: Presentation, from_idx: int, to_idx: int) -> None:
    sldIdLst = prs.slides._sldIdLst
    slides   = list(sldIdLst)
    if not (0 <= from_idx < len(slides)):
        return
    to_idx = max(0, min(to_idx, len(slides) - 1))
    sldId  = slides[from_idx]
    sldIdLst.remove(sldId)
    sldIdLst.insert(to_idx, sldId)


# ── Profit pool year resolution ───────────────────────────────────────────────

def _latest_available_profit_pool_year(
    df: pd.DataFrame,
    requested_year: Optional[int] = None,
) -> int:
    available: List[int] = []
    for col in df.columns:
        m = _re.search(r"Total Revenue \[FY (\d{4})\]", str(col))
        if m:
            y = int(m.group(1))
            if f"EBITDA [FY {y}] ($USDmm, Historical rate)" in df.columns:
                available.append(y)
    if not available:
        return DEFAULT_PROFIT_POOL_YEAR
    cap        = int(requested_year) if requested_year else max(available)
    candidates = [y for y in available if y <= cap]
    return max(candidates) if candidates else min(available)


# ── Growth slide data builder ─────────────────────────────────────────────────

def _growth_slide_to_tc_items(
    content: da.GrowthSlideContent,
    title_text: str,
    chart_elem_name: str = TC_ELEM_CHART,
    title_elem_name: str = TC_ELEM_TITLE,
    start_year: int = GROWTH_START_YEAR,
    end_year:   int = GROWTH_END_YEAR,
    hist_start: int = HIST_START,
    hist_end:   int = HIST_END,
    fc_start:   int = FC_START,
    fc_end:     int = FC_END,
) -> List[dict]:
    yoy_by_year = {
        int(y): v
        for y, v in zip(content.years or [], content.yoy or [])
        if y is not None
    }
    slide_years = list(range(start_year, end_year + 1))

    def scaled(v: Optional[float]) -> Optional[Dict[str, float]]:
        return _tc_number(v * 100.0) if v is not None else None

    hist_row: List[Optional[Dict]] = []
    fc_row:   List[Optional[Dict]] = []

    for y in slide_years:
        val = yoy_by_year.get(y)
        if hist_start <= y <= hist_end:
            hist_row.append(scaled(val))
        elif y == fc_start:
            hist_row.append(scaled(yoy_by_year.get(fc_start)))
        else:
            hist_row.append(None)
        fc_row.append(scaled(val) if fc_start <= y <= fc_end else None)

    table = [
        [_tc_string("")]                                          + [_tc_year(y)   for y in slide_years],
        [_tc_string(f"Historical ({hist_start}-{hist_end})")]    + hist_row,
        [_tc_string(f"Forecast ({fc_start}-{fc_end})")]          + fc_row,
    ]
    return [
        {"name": chart_elem_name, "table": table},
        {"name": title_elem_name, "table": [[_tc_string(title_text)]]},
    ]


# ── Profit pool data builder ──────────────────────────────────────────────────

def _revenue_col(year: int) -> str:
    return f"Total Revenue [FY {year}] ($USDmm, Historical rate)"


def _ebitda_col(year: int) -> str:
    return f"EBITDA [FY {year}] ($USDmm, Historical rate)"


def _normalize_tag_series(s: pd.Series) -> pd.Series:
    return (
        s.fillna("").astype(str)
        .str.replace("\u00A0", " ", regex=False)
        .str.strip()
        .str.replace(r"\s+", " ", regex=True)
        .str.title()
    )


def _clean_str(s: pd.Series) -> pd.Series:
    return s.fillna("").astype(str).map(str.strip)


def _build_profit_pool_view(
    df: pd.DataFrame,
    year: int,
    country_filter: Optional[str] = None,
) -> pd.DataFrame:
    rev_col  = _revenue_col(year)
    ebt_col  = _ebitda_col(year)
    required = [COL_FINAL_TAG, COL_SIC, rev_col, ebt_col, COL_COUNTRY]
    missing  = [c for c in required if c not in df.columns]
    if missing:
        raise ValueError(f"Missing columns for Profit Pool slide: {missing}")

    d = df.copy()
    if country_filter and country_filter.strip().lower() not in ("", "global"):
        d = d[_clean_str(d[COL_COUNTRY]) == country_filter.strip()].copy()

    d[COL_FINAL_TAG] = _normalize_tag_series(d[COL_FINAL_TAG])
    d[COL_SIC]       = _clean_str(d[COL_SIC])
    d[rev_col]       = pd.to_numeric(d[rev_col], errors="coerce").fillna(0.0)
    d[ebt_col]       = pd.to_numeric(d[ebt_col], errors="coerce").fillna(0.0)

    cement_mask = (
        (d[COL_FINAL_TAG].str.lower() == CEMENT_PARENT.lower())
        & d[COL_SIC].str.contains(r"cement", case=False, na=False)
    )
    cement_rev = float(d.loc[cement_mask, rev_col].sum())
    cement_ebt = float(d.loc[cement_mask, ebt_col].sum())

    rows = []
    for cat in [c for c in d[COL_FINAL_TAG].unique() if c]:
        mask = (
            (d[COL_FINAL_TAG].str.lower() == CEMENT_PARENT.lower()) & ~cement_mask
            if cat.lower() == CEMENT_PARENT.lower()
            else d[COL_FINAL_TAG] == cat
        )
        rev = float(d.loc[mask, rev_col].sum())
        ebt = float(d.loc[mask, ebt_col].sum())
        rows.append({"Category": cat, "Revenue": rev, "EBITDA": ebt,
                     "EBITDA margin": ebt / rev if rev else 0.0})

    rows.append({"Category": "Cement", "Revenue": cement_rev, "EBITDA": cement_ebt,
                 "EBITDA margin": cement_ebt / cement_rev if cement_rev else 0.0})

    df_view = (
        pd.DataFrame(rows)
        .loc[lambda x: x["Category"].str.strip() != ""]
        .assign(Category=lambda x: x["Category"].replace(CATEGORY_RENAMES))
        .sort_values("EBITDA margin", ascending=False)
        .reset_index(drop=True)
    )

    total_rev = float(df_view["Revenue"].sum())
    if total_rev <= 0:
        raise ValueError("Total revenue is 0; cannot build Profit Pool slide.")

    df_view["width"] = df_view["Revenue"].clip(lower=0.0) / total_rev
    return df_view[["Category", "Revenue", "EBITDA", "EBITDA margin", "width"]].copy()


def _profit_pool_to_tc_items(df_flat: pd.DataFrame, year: int, country: str) -> List[dict]:
    df_view = _build_profit_pool_view(df_flat, year=year, country_filter=country)
    cats    = df_view["Category"].astype(str).tolist()
    margins = df_view["EBITDA margin"].astype(float).tolist()
    widths  = df_view["Revenue"].astype(float).tolist()

    table = [
        [_tc_string("")]              + [_tc_string(c)        for c in cats],
        [_tc_string("Revenue")]       + [_tc_number(v)        for v in widths],
        [_tc_string("EBITDA margin")] + [_tc_percent(v * 100) for v in margins],
    ]
    title = (
        f"{country} \u2014 EBITDA margin by category (FY {year})"
        if country else f"EBITDA margin by category (FY {year})"
    )
    return [
        {"name": TC_ELEM_CATEGORY_CHART, "table": table},
        {"name": TC_ELEM_CATEGORY_TITLE, "table": [[_tc_string(title)]]},
    ]


def _build_subtitle(country: Optional[str], company: Optional[str]) -> Optional[str]:
    if country and company:
        return f"{company} in {country}"
    return country or company or None


# ── Main build function ───────────────────────────────────────────────────────

def build_deck(req: DeckRequest) -> Presentation:
    """
    Assemble a multi-slide PPTX deck exactly matching the Streamlit builder:
      1  Executive Summary (python-pptx, prepended)
      2  Construction YoY growth (think-cell)
      3  Building products growth (think-cell)
      4  Cement growth (think-cell)
      5  Profit pool (think-cell, optional)
      6  KPI Index (python-pptx, only when KPI slides present)
      7+ KPI comparison slides (think-cell)
    """
    theme    = default_theme()
    country  = (req.country or "").strip() or "Global"
    subtitle = _build_subtitle(req.country, req.company)

    # Resolve template paths from settings if not provided in request
    tc_url = req.tc_server_url or str(settings.THINKCELL_SERVER_URL)
    tmpl_s2 = req.tc_growth_template_s2 or str(settings.THINKCELL_TEMPLATE_GROWTH_S2)
    tmpl_s3 = req.tc_growth_template_s3 or str(settings.THINKCELL_TEMPLATE_GROWTH_S3)
    tmpl_s4 = req.tc_growth_template_s4 or str(settings.THINKCELL_TEMPLATE_GROWTH_S4)
    tmpl_cat = req.tc_category_template or str(settings.THINKCELL_TEMPLATE_CATEGORY)

    # ── Get exec summary content ──────────────────────────────────────────────
    exec_content = da.get_exec_summary_content(
        country=country,
        company=req.company,
        slide_polish=True,
        use_web_reasons=True,
        use_cache=True,
    )

    # ── Fallback: no flat file data ────────────────────────────────────────────
    print(f"[DECK] df_flat is None: {req.df_flat is None}", flush=True)
    print(f"[DECK] df_flat empty: {req.df_flat.empty if req.df_flat is not None else 'N/A'}", flush=True)
    print(f"[DECK] df_flat shape: {req.df_flat.shape if req.df_flat is not None else 'N/A'}", flush=True)
    print(f"[DECK] df_profit_pool shape: {req.df_profit_pool.shape if req.df_profit_pool is not None else 'N/A'}", flush=True)
    print(f"[DECK] country={country}, company={req.company}, year={req.year}", flush=True)
    print(f"[DECK] tc_url={tc_url}", flush=True)
    print(f"[DECK] tmpl_s2={tmpl_s2}", flush=True)
    if req.df_flat is None or req.df_flat.empty:
        print("[DECK] WARNING: df_flat is empty — returning single exec summary slide", flush=True)
        prs = Presentation()
        prs.slide_width  = theme.slide_width
        prs.slide_height = theme.slide_height
        add_slide_01_exec_summary(prs, theme, exec_content, subtitle=subtitle)
        return prs

    # ── Slides 2–4: growth series via think-cell ──────────────────────────────
    c2 = da.compute_growth_view_series(req.df_flat, category="Construction overall",                country=country)
    c3 = da.compute_growth_view_series(req.df_flat, category="Building Products Overall Sales",     country=country)
    c4 = da.compute_growth_view_series(req.df_flat, category="Cement, Concrete, Lime Overall Sales", country=country)

    tc_items: List[Dict[str, Any]] = [
        {
            "template": tmpl_s2,
            "data": _growth_slide_to_tc_items(c2, title_text=f"{country} \u2014 Construction overall YoY growth"),
        },
        {
            "template": tmpl_s3,
            "data": _growth_slide_to_tc_items(c3, title_text=f"{country} \u2014 Building products YoY growth"),
        },
        {
            "template": tmpl_s4,
            "data": _growth_slide_to_tc_items(c4, title_text=f"{country} \u2014 Cement / concrete / lime YoY growth"),
        },
    ]

    # ── Slide 5: Profit pool ──────────────────────────────────────────────────
    has_profit_pool = req.df_profit_pool is not None and not req.df_profit_pool.empty
    if has_profit_pool:
        try:
            profit_pool_year = _latest_available_profit_pool_year(
                req.df_profit_pool, requested_year=req.year
            )
            tc_items.append({
                "template": tmpl_cat,
                "data": _profit_pool_to_tc_items(req.df_profit_pool, year=profit_pool_year, country=country),
            })
        except Exception as e:
            print(f"[DECK] Profit pool slide skipped: {e}", flush=True)
            has_profit_pool = False

    # ── Slides 7+: KPI comparison slides ─────────────────────────────────────
    has_kpi = (
        req.comparison_request is not None
        and bool(req.comparison_request.kpi_selections)
    )
    print(f"[DECK] has_kpi={has_kpi}", flush=True)
    if has_kpi:
        cr = req.comparison_request
        print(f"[DECK] KPI base={cr.base_company}, peers={cr.peer_companies}", flush=True)
        print(f"[DECK] KPI selections={[(s.kpi_key, s.chart_mode) for s in cr.kpi_selections]}", flush=True)
        print(f"[DECK] wide_df={cr.wide_df.shape if cr.wide_df is not None else None}", flush=True)
        print(f"[DECK] long_df={cr.long_df.shape if cr.long_df is not None else None}", flush=True)
        print(f"[DECK] bar_template={cr.bar_template_path}", flush=True)
        print(f"[DECK] line_template={cr.line_template_path}", flush=True)
    slide_06_payloads: List[Dict[str, Any]] = []
    if has_kpi:
        slide_06_payloads = build_slide_06_tc_payloads(req.comparison_request)
        print(f"[DECK] slide_06_payloads count={len(slide_06_payloads)}", flush=True)

    # ── Single think-cell POST for all chart slides ───────────────────────────
    all_tc_items = tc_items + slide_06_payloads
    print(f"[DECK] Sending {len(all_tc_items)} items to think-cell", flush=True)
    for i, item in enumerate(all_tc_items):
        print(f"[DECK]   item[{i}] template={item['template']}", flush=True)
    ppttc_json   = _build_multi_ppttc_payload(all_tc_items)
    print(f"[DECK] ppttc payload length: {len(ppttc_json)} chars", flush=True)
    pptx_bytes   = _post_to_tcserver(tc_url, ppttc_json)
    print(f"[DECK] think-cell returned {len(pptx_bytes)} bytes", flush=True)

    prs = Presentation(io.BytesIO(pptx_bytes))
    print(f"[DECK] Slides from think-cell: {len(prs.slides)}", flush=True)
    _fix_axis_font_sizes(prs, target_pt=10.0)
    prs.slide_width  = theme.slide_width
    prs.slide_height = theme.slide_height

    # ── Prepend Slide 1: Executive Summary ────────────────────────────────────
    add_slide_01_exec_summary(prs, theme, exec_content, subtitle=subtitle)
    print(f"[DECK] Slides after exec summary: {len(prs.slides)}", flush=True)
    _move_slide(prs, from_idx=len(prs.slides) - 1, to_idx=0)
    print(f"[DECK] Final slide count: {len(prs.slides)}", flush=True)

    # ── Insert Slide 6: KPI Index (after base slides, before KPI charts) ──────
    if has_kpi and slide_06_payloads:
        n_base_slides        = 5 if has_profit_pool else 4
        first_kpi_slide_num  = n_base_slides + 2  # +1 exec, +1 for index itself

        build_kpi_index_slide(
            prs,
            req.comparison_request,
            first_kpi_slide_number=first_kpi_slide_num,
        )
        index_target = n_base_slides  # 0-based: after slide 1..n_base
        _move_slide(prs, from_idx=len(prs.slides) - 1, to_idx=index_target)

    return prs